import warnings
warnings.filterwarnings("ignore", message="Thread 'MainThread': missing ScriptRunContext!")
warnings.filterwarnings("ignore", message="No runtime found, using MemoryCacheStorageManager")

import streamlit as st
import pandas as pd
import numpy as np
from datetime import datetime
import requests
from io import BytesIO

try:
    import matplotlib.pyplot as plt
    MATPLOTLIB_AVAILABLE = True
except ImportError:
    MATPLOTLIB_AVAILABLE = False
    st.warning("Matplotlib не установлен. Некоторые графики будут недоступны.")

try:
    import seaborn as sns
    SEABORN_AVAILABLE = True
except ImportError:
    SEABORN_AVAILABLE = False

try:
    from sklearn.model_selection import train_test_split
    from sklearn.preprocessing import PolynomialFeatures, StandardScaler, OneHotEncoder
    from sklearn.linear_model import LinearRegression
    from sklearn.pipeline import Pipeline
    from sklearn.compose import ColumnTransformer
    from sklearn.impute import SimpleImputer
    from sklearn.metrics import mean_squared_error, r2_score, mean_absolute_error
    SKLEARN_AVAILABLE = True
except ImportError:
    SKLEARN_AVAILABLE = False
    st.error("Библиотека scikit-learn не установлена. Функция прогнозирования будет недоступна.")

try:
    import tempfile
    import os
    import base64
    from pptx import Presentation
    from pptx.util import Inches
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False

st.set_page_config(
    page_title="Анализ рынка недвижимости - Полный анализ",
    page_icon="🏠",
    layout="wide",
    initial_sidebar_state="expanded"
)

@st.cache_data(ttl=3600)  
def load_data_from_github():
    github_url = "https://github.com/evoluteorburn-lab/HSE_exam_DS16_GusevA/raw/5a932eac450c1cc46fd032db7deb0fc9dd86843b/Cian.xlsx"
    
    try:
        response = requests.get(github_url)
        response.raise_for_status()
        file_content = BytesIO(response.content)
        df = pd.read_excel(file_content)
        if 'Номер квартиры' not in df.columns:
            df['Номер квартиры'] = range(1, len(df) + 1)
        return df
    except Exception as e:
        st.error(f"Ошибка при загрузке данных: {e}")
        return create_demo_data()

def create_demo_data():
    data = {
        'ID Корпуса': [1, 2, 3, 4, 5, 6, 7, 8, 9, 10],
        'ID ЖК': [101, 101, 102, 103, 103, 104, 105, 106, 107, 108],
        'ЖК рус': ['ЖК А', 'ЖК А', 'ЖК Б', 'ЖК В', 'ЖК В', 'ЖК Г', 'ЖК Д', 'ЖК Е', 'ЖК Ж', 'ЖК З'],
        'Комнат': [1, 2, 3, 1, 2, 3, 4, 1, 2, 3],
        'Площадь': [30.5, 45.2, 60.1, 25.0, 40.0, 75.3, 90.0, 28.5, 42.0, 65.0],
        'Цена': [5000000, 8000000, 12000000, 4000000, 7000000, 15000000, 20000000, 4500000, 7500000, 13000000],
        'Цена кв м': [163934, 176991, 199667, 160000, 175000, 199203, 222222, 157895, 178571, 200000],
        'Этаж': [5, 8, 3, 2, 7, 12, 15, 4, 6, 9],
        'Район Город': ['ЦАО', 'САО', 'ЮАО', 'ЦАО', 'САО', 'ЮАО', 'ЗАО', 'СВАО', 'ЮЗАО', 'ВАО'],
        'Тип помещения': ['Квартира', 'Квартира', 'Апартаменты', 'Квартира', 'Апартаменты', 'Квартира', 'Квартира', 'Квартира', 'Апартаменты', 'Квартира'],
        'Застройщик': ['ПИК', 'Самолет', 'Эталон', 'ПИК', 'Самолет', 'Эталон', 'ПИК', 'Самолет', 'Эталон', 'ПИК'],
        'Класс К....': ['Комфорт', 'Бизнес', 'Премиум', 'Эконом', 'Комфорт', 'Бизнес', 'Премиум', 'Эконом', 'Комфорт', 'Бизнес'],
        'Отделка помещения': ['Да', 'Нет', 'Да', 'Нет', 'Да', 'Да', 'Нет', 'Да', 'Нет', 'Да'],
        'Школа/Детский Сад': ['Рядом', 'Далеко', 'Рядом', 'Далеко', 'Рядом', 'Рядом', 'Далеко', 'Рядом', 'Далеко', 'Рядом'],
        'Парк/Зона отдыха': ['Рядом', 'Далеко', 'Рядом', 'Далеко', 'Рядом', 'Рядом', 'Далеко', 'Рядом', 'Далеко', 'Рядом'],
        'Спорт': ['Есть', 'Нет', 'Есть', 'Нет', 'Есть', 'Есть', 'Нет', 'Есть', 'Нет', 'Есть'],
        'Парковка': ['Подземная', 'Улица', 'Подземная', 'Улица', 'Подземная', 'Подземная', 'Улица', 'Подземная', 'Улица', 'Подземная'],
        'Рестораны': ['Много', 'Мало', 'Много', 'Мало', 'Много', 'Много', 'Мало', 'Много', 'Мало', 'Много'],
        'Номер квартиры': [1, 2, 3, 4, 5, 6, 7, 8, 9, 10]
    }
    return pd.DataFrame(data)

def create_presentation(filtered_data, analysis_results=None):
    if not PPTX_AVAILABLE:
        st.error("Библиотека python-pptx не установлена. Невозможно создать презентацию.")
        return None
    
    prs = Presentation()
    
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "Анализ рынка недвижимости"
    subtitle.text = f"Отчет сгенерирован {datetime.now().strftime('%d.%m.%Y %H:%M')}\nНайдено объектов: {len(filtered_data)}"
    
    stats_slide_layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(stats_slide_layout)
    title = slide.shapes.title
    title.text = "Общая статистика"
    
    price_column = 'Цена кв м' if 'Цена кв м' in filtered_data.columns else 'Цена'
    if price_column in filtered_data.columns:
        avg_price = filtered_data[price_column].mean()
        max_price = filtered_data[price_column].max()
        min_price = filtered_data[price_column].min()
        
        rows = [
            ["Показатель", "Значение"],
            ["Количество объектов", str(len(filtered_data))],
            ["Максимальная цена", f"{max_price:,.0f} руб."],
            ["Средняя цена", f"{avg_price:,.0f} руб."],
            ["Минимальная цена", f"{min_price:,.0f} руб."]
        ]
        
        x, y, cx, cy = Inches(1), Inches(1.5), Inches(8), Inches(2)
        shape = slide.shapes.add_table(len(rows), 2, x, y, cx, cy)
        table = shape.table
        
        for i, row in enumerate(rows):
            for j, cell_value in enumerate(row):
                table.cell(i, j).text = cell_value
    
    if analysis_results:
        analysis_slide_layout = prs.slide_layouts[5]
        slide = prs.slides.add_slide(analysis_slide_layout)
        title = slide.shapes.title
        title.text = "Результаты прогнозирования"
        
        content = f"""
        Модель прогнозирования обучена на {analysis_results.get('train_samples', 0)} объектах
        Качество модели (R²): {analysis_results.get('r2_score', 0):.3f}
        Средняя ошибка прогноза (RMSE): {analysis_results.get('rmse', 0):.2f}
        """
        
        txBox = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(2))
        tf = txBox.text_frame
        tf.text = content
    
    temp_file = tempfile.NamedTemporaryFile(delete=False, suffix='.pptx')
    prs.save(temp_file.name)
    
    return temp_file.name

data = load_data_from_github()

st.title("🏠 Комплексный анализ рынка недвижимости")
st.write(f"Данные загружены: {len(data)} строк, {len(data.columns)} колонок")

if not SKLEARN_AVAILABLE:
    st.error("⚠️ Библиотека scikit-learn не установлена. Функция прогнозирования будет недоступна.")

if not MATPLOTLIB_AVAILABLE:
    st.warning("⚠️ Библиотека matplotlib не установлена. Некоторые графики будут недоступны.")

st.sidebar.title("Навигация")
section_options = ["Поиск квартир"]
if SKLEARN_AVAILABLE:
    section_options.append("Прогнозирование")

section = st.sidebar.radio("Выберите раздел:", section_options)

if 'filtered_apartments' not in st.session_state:
    st.session_state.filtered_apartments = None
if 'trained_model' not in st.session_state:
    st.session_state.trained_model = None
if 'model_features' not in st.session_state:
    st.session_state.model_features = None
if 'target_column' not in st.session_state:
    st.session_state.target_column = None
if 'analysis_results' not in st.session_state:
    st.session_state.analysis_results = None

def show_apartment_search():
    st.header("🔍 Поиск квартир по параметрам")
    
    col1, col2 = st.columns(2)
    
    with col1:
        st.subheader("Основные параметры")
        
        def get_unique_values(column_name, default_options=None):
            if column_name in data.columns and not data[column_name].empty:
                unique_vals = data[column_name].dropna().unique().tolist()
                return sorted([x for x in unique_vals if x is not None and x != ''])
            return default_options if default_options else []
        
        class_options = get_unique_values('Класс К....', ['Эконом', 'Комфорт', 'Бизнес', 'Премиум'])
        class_input = st.selectbox('Класс квартиры', options=[None] + class_options)
        
        area_min = st.number_input('Площадь от (м²)', min_value=0.0, value=0.0)
        area_max = st.number_input('Площадь до (m²)', min_value=0.0, value=0.0)
    
    with col2:
        st.subheader("Дополнительные параметры")
        
        rooms_options = get_unique_values('Комнат', [1, 2, 3, 4, 5])
        rooms_input = st.selectbox('Комнат', options=[None] + rooms_options)
        
        floor_min = st.number_input('Этаж от', min_value=0, value=0)
        floor_max = st.number_input('Этаж до', min_value=0, value=0)
        
        district_options = get_unique_values('Район Город', ['ЦАО', 'САО', 'ЮАО', 'ЗАО', 'СВАО', 'ЮЗАО', 'ВАО'])
        district_input = st.selectbox('Район', options=[None] + district_options)
        
        builder_options = get_unique_values('Застройщик', ['ПИК', 'Самолет', 'Эталон'])
        builder_input = st.selectbox('Застройщик', options=[None] + builder_options)
    
    st.subheader("🏗️ Инфраструктура")
    infra_cols = st.columns(5)
    infrastructure_options = {}
    
    infra_columns = ['Школа/Детский Сад', 'Парк/Зона отдыха', 'Спорт', 'Парковка', 'Рестораны']
    
    for i, col_name in enumerate(infra_columns):
        if col_name in data.columns:
            options = get_unique_values(col_name, [])
            with infra_cols[i]:
                infrastructure_options[col_name] = st.selectbox(
                    col_name,
                    options=[None] + options,
                    key=f"infra_{col_name}"
                )
    
    if st.button('Найти квартиры', type='primary'):
        filtered_df = data.copy()
        
        if class_input and 'Класс К....' in filtered_df.columns:
            filtered_df = filtered_df[filtered_df['Класс К....'] == class_input]
        
        if area_min > 0:
            filtered_df = filtered_df[filtered_df['Площадь'] >= area_min]
        if area_max > 0:
            filtered_df = filtered_df[filtered_df['Площадь'] <= area_max]
        
        if rooms_input:
            filtered_df = filtered_df[filtered_df['Комнат'] == rooms_input]
        
        if floor_min > 0:
            filtered_df = filtered_df[filtered_df['Этаж'] >= floor_min]
        if floor_max > 0:
            filtered_df = filtered_df[filtered_df['Этаж'] <= floor_max]
        
        if district_input:
            filtered_df = filtered_df[filtered_df['Район Город'] == district_input]
        
        if builder_input:
            filtered_df = filtered_df[filtered_df['Застройщик'] == builder_input]
        
        for col_name, value in infrastructure_options.items():
            if value and col_name in filtered_df.columns:
                filtered_df = filtered_df[filtered_df[col_name] == value]
        
        if len(filtered_df) == 0:
            st.warning("Не найдено объектов по указанным параметрам")
            st.session_state.filtered_apartments = None
        else:
            st.success(f"Найдено {len(filtered_df)} объектов")
            st.session_state.filtered_apartments = filtered_df
            
            price_column = 'Цена кв м' if 'Цена кв м' in filtered_df.columns else 'Цена'
            if price_column in filtered_df.columns:
                col1, col2, col3 = st.columns(3)
                with col1:
                    st.metric("Максимальная цена", f"{filtered_df[price_column].max():,.0f} руб.")
                with col2:
                    st.metric("Средняя цена", f"{filtered_df[price_column].mean():,.0f} руб.")
                with col3:
                    st.metric("Минимальная цена", f"{filtered_df[price_column].min():,.0f} руб.")
            
            display_columns = ['Номер квартиры', 'Площадь', 'Комнат', 'Этаж', 'Район Город', 'Цена кв м', 'Класс К....']
            display_columns.extend([col for col in infra_columns if col in filtered_df.columns])
            
            available_columns = [col for col in display_columns if col in filtered_df.columns]
            
            display_df = filtered_df[available_columns].copy()
            display_df.rename(columns={
                'Класс К....': 'Класс',
                'Район Город': 'Район',
                'Цена кв м': 'Цена за м²',
                'Номер квартиры': '№ Квартиры'
            }, inplace=True)
            
            st.dataframe(
                display_df.style.format({
                    'Цена за м²': '{:,.0f} руб.',
                    'Площадь': '{:.1f} м²',
                    '№ Квартиры': '{:.0f}'
                }),
                height=400
            )
            
            if PPTX_AVAILABLE:
                if st.button("📊 Выгрузить презентацию по результатам анализа"):
                    presentation_path = create_presentation(filtered_df)
                    if presentation_path:
                        with open(presentation_path, "rb") as file:
                            btn = st.download_button(
                                label="⬇️ Скачать презентацию",
                                data=file,
                                file_name=f"анализ_недвижимости_{datetime.now().strftime('%Y%m%d_%H%M')}.pptx",
                                mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
                            )
                        os.unlink(presentation_path)
            else:
                st.warning("Функция создания презентации недоступна (требуется python-pptx)")
            
            if SKLEARN_AVAILABLE:
                if st.button("📊 Сделать прогноз для найденных квартир"):
                    st.session_state.target_section = "Прогнозирование"
                    st.rerun()
            else:
                st.warning("Функция прогнозирования недоступна (требуется scikit-learn)")

def show_polynomial_regression():
    if not SKLEARN_AVAILABLE:
        st.error("Функция прогнозирования недоступна. Установите scikit-learn.")
        return
        
    st.header("📈 Прогнозирование цен на 2026 год")
    
    st.info("""
    В этом разделе строится полиномиальная регрессионная модель для прогнозирования цены за м² 
    на основе выбранных признаков. Модель может быть обучена на всех данных или на отфильтрованных квартирах.
    """)
    
    use_filtered = st.checkbox("Использовать отфильтрованные квартиры из поиска", 
                              value=st.session_state.filtered_apartments is not None,
                              disabled=st.session_state.filtered_apartments is None)
    
    if use_filtered and st.session_state.filtered_apartments is not None:
        analysis_data = st.session_state.filtered_apartments
        st.info(f"Используются отфильтрованные данные: {len(analysis_data)} квартир")
    else:
        analysis_data = data
        st.info("Используются все данные")
    
    numeric_cols = analysis_data.select_dtypes(include=['int64', 'float64']).columns.tolist()
    categorical_cols = analysis_data.select_dtypes(include=['object']).columns.tolist()
    
    target_col = st.selectbox("Целевая переменная (y):", 
                             options=['Цена кв м', 'Цена'],
                             index=0)
    
    if target_col not in analysis_data.columns:
        st.error(f"Колонка '{target_col}' не найдена в данных!")
        return
    
    price_columns_to_exclude = ['Цена кв м', 'Цена', 'Цена со скидкой', 'Изменение цены последнее', 'Изменение цены', 'Номер квартиры']
    available_features = [col for col in numeric_cols + categorical_cols 
                         if col not in price_columns_to_exclude]
    
    available_numeric = [col for col in available_features if col in numeric_cols]
    available_categorical = [col for col in available_features if col in categorical_cols]
    
    col1, col2 = st.columns([4, 1])
    
    with col1:
        selected_features = st.multiselect("Признаки для модели (X):", 
                                          options=available_features,
                                          default=['Площадь', 'Комнат', 'Этаж'])
    
    with col2:
        st.write("")  # Отступ для выравнивания
        st.write("")  # Отступ для выравнивания
        
        if st.button("Все признаки", key="select_all_btn"):
            selected_features = available_features
            st.rerun()
        
        if st.button("Все категориальные", key="select_categorical_btn"):
            # Добавляем категориальные признаки, сохраняя уже выбранные
            current_selected = set(selected_features)
            current_selected.update(available_categorical)
            selected_features = list(current_selected)
            st.rerun()
            
        if st.button("Все числовые", key="select_numeric_btn"):
            # Добавляем числовые признаки, сохраняя уже выбранные
            current_selected = set(selected_features)
            current_selected.update(available_numeric)
            selected_features = list(current_selected)
            st.rerun()

        if st.button("Очистить", key="clear_btn"):
            selected_features = []
            st.rerun()
    
    if not selected_features:
        st.warning("Выберите хотя бы один признак для построения модели")
        return

    selected_numeric = [col for col in selected_features if col in numeric_cols]
    selected_categorical = [col for col in selected_features if col in categorical_cols]
    
    st.info(f"Выбрано: {len(selected_features)} признаков "
           f"({len(selected_numeric)} числовых, {len(selected_categorical)} категориальных)")
    
    col1, col2, col3 = st.columns(3)
    with col1:
        degree = st.slider("Степень полинома", 1, 5, 2)
    with col2:
        test_size = st.slider("Размер тестовой выборки", 0.1, 0.5, 0.3)
    with col3:
        random_state = st.number_input("Random state", 0, 100, 42)
    
    if st.button("Обучить модель", type='primary'):
        try:
            
            missing_features = [col for col in selected_features if col not in analysis_data.columns]
            if missing_features:
                st.error(f"Следующие признаки не найдены в данных: {missing_features}")
                return
            
            X = analysis_data[selected_features].copy()
            y = analysis_data[target_col]

            X_clean = X.dropna()
            y_clean = y.loc[X_clean.index]
            
            if len(X_clean) == 0:
                st.error("После обработки пропусков не осталось данных для обучения!")
                return
            
            if len(X_clean) < 10:
                st.warning(f"Мало данных для обучения: всего {len(X_clean)} строк. Результаты могут быть ненадежными.")
            
            numeric_features = X_clean.select_dtypes(include=['int64', 'float64']).columns.tolist()
            categorical_features = X_clean.select_dtypes(include=['object']).columns.tolist()

            if not numeric_features and not categorical_features:
                st.error("Не осталось признаков для обучения после обработки данных!")
                return
            
            numeric_transformer = Pipeline(steps=[
                ('imputer', SimpleImputer(strategy='median')),
                ('scaler', StandardScaler())
            ])
            
            categorical_transformer = Pipeline(steps=[
                ('imputer', SimpleImputer(strategy='most_frequent')),
                ('onehot', OneHotEncoder(handle_unknown='ignore', sparse_output=False))
            ])
            
            transformers = []
            if numeric_features:
                transformers.append(('num', numeric_transformer, numeric_features))
            if categorical_features:
                transformers.append(('cat', categorical_transformer, categorical_features))
            
            preprocessor = ColumnTransformer(transformers=transformers)
            
            full_pipeline = Pipeline(steps=[
                ('preprocessor', preprocessor),
                ('poly', PolynomialFeatures(degree=degree, include_bias=False))
            ])
            
            X_processed = full_pipeline.fit_transform(X_clean)

            if X_processed.shape[0] == 0:
                st.error("После преобразований не осталось данных для обучения!")
                return
            
            X_train, X_test, y_train, y_test = train_test_split(
                X_processed, y_clean, test_size=test_size, random_state=random_state
            )
            
            if len(X_test) == 0:
                st.error("Тестовая выборка пустая! Уменьшите размер тестовой выборки.")
                return
            
            model = LinearRegression()
            model.fit(X_train, y_train)
            
            y_pred = model.predict(X_test)
            
            mae = mean_absolute_error(y_test, y_pred)
            mse = mean_squared_error(y_test, y_pred)
            rmse = np.sqrt(mse)
            r2 = r2_score(y_test, y_pred)
            
            st.success("Модель успешно обучена!")
            
            st.session_state.trained_model = model
            st.session_state.model_pipeline = full_pipeline
            st.session_state.model_features = selected_features
            st.session_state.target_column = target_col
            
            st.session_state.analysis_results = {
                'train_samples': len(X_train),
                'r2_score': r2,
                'rmse': rmse,
                'mae': mae
            }
            
            try:
                all_predictions = model.predict(full_pipeline.transform(X))
                
                forecast_df = analysis_data[['Номер квартиры'] + selected_features].copy()
                forecast_df['Фактическая цена'] = analysis_data[target_col]
                forecast_df['Прогноз на 2026 год'] = all_predictions
                forecast_df['Изменение, %'] = ((forecast_df['Прогноз на 2026 год'] - forecast_df['Фактическая цена']) / 
                                             forecast_df['Фактическая цена'] * 100)
                
                st.subheader("📊 Прогноз цен на 2026 год")
                st.dataframe(
                    forecast_df.style.format({
                        'Фактическая цена': '{:,.0f}',
                        'Прогноз на 2026 год': '{:,.0f}',
                        'Изменение, %': '{:.1f}%'
                    }),
                    height=400
                )
                
            except Exception as e:
                st.warning(f"Не удалось сделать прогноз для всех данных: {e}")
                st.info("Отображаются результаты только для тестовой выборки")
            
            col1, col2, col3, col4 = st.columns(4)
            with col1:
                st.metric("R² Score", f"{r2:.3f}")
            with col2:
                st.metric("RMSE", f"{rmse:.2f}")
            with col3:
                st.metric("MAE", f"{mae:.2f}")
            with col4:
                st.metric("Обучено на", f"{len(X_train)} samples")
            
            test_results_df = pd.DataFrame({
                'Actual': y_test.values,
                'Predicted': y_pred,
                'Residual': y_test.values - y_pred
            }).round(2)
            
            with st.expander("Детали предсказаний (тестовая выборка)"):
                st.dataframe(test_results_df.head(10))
            
            if PPTX_AVAILABLE and st.session_state.filtered_apartments is not None:
                if st.button("📊 Выгрузить презентацию с результатами прогноза"):
                    presentation_path = create_presentation(
                        st.session_state.filtered_apartments, 
                        st.session_state.analysis_results
                    )
                    if presentation_path:
                        with open(presentation_path, "rb") as file:
                            btn = st.download_button(
                                label="⬇️ Скачать презентацию с прогнозом",
                                data=file,
                                file_name=f"прогноз_недвижимости_{datetime.now().strftime('%Y%m%d_%H%M')}.pptx",
                                mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
                            )
                        os.unlink(presentation_path)
            
        except Exception as e:
            st.error(f"Ошибка при обучении модели: {str(e)}")
            st.info("Попробуйте выбрать меньше признаков или уменьшить степень полинома")

if 'target_section' in st.session_state:
    section = st.session_state.target_section
    del st.session_state.target_section

if section == "Поиск квартир":
    show_apartment_search()
elif section == "Прогнозирование" and SKLEARN_AVAILABLE:
    show_polynomial_regression()
